"""Build the Saab internal briefing deck."""

from __future__ import annotations

from pathlib import Path

import pymupdf
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
THESIS = ROOT.parent
FIG = THESIS / "figures" / "experiments"
GEN = ROOT / "assets" / "generated"
CONV = ROOT / "assets" / "converted"
ANIM = ROOT / "assets" / "animations"
CONV.mkdir(parents=True, exist_ok=True)

W, H = Inches(13.333), Inches(7.5)
BG = RGBColor(255, 255, 255)
INK = RGBColor(20, 32, 52)
MUTED = RGBColor(90, 104, 122)
CYAN = RGBColor(10, 122, 150)
ORANGE = RGBColor(196, 90, 32)
WHITE = RGBColor(255, 255, 255)
NAVY = RGBColor(20, 32, 52)
CARD = RGBColor(243, 246, 248)
LINE = RGBColor(210, 218, 226)
FOOTER = RGBColor(245, 247, 249)


def _set_run(run, size, color, bold=False, name="Calibri"):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = name
    run.font.italic = False


def _fill(shape, color, line=None, line_width=1.0):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_width)


def _textbox(slide, left, top, width, height, text, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf._txBody.bodyPr.set("anchor", {MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"}[anchor])
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    _set_run(run, size, color, bold=bold)
    return box


def _bullets(slide, left, top, width, items, size=18, color=INK, gap=0.42):
    y = top
    for item in items:
        mark = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, y + Inches(0.10), Inches(0.12), Inches(0.12))
        _fill(mark, CYAN)
        _textbox(slide, left + Inches(0.28), y, width - Inches(0.28), Inches(gap + 0.08), item, size=size, color=color)
        y += Inches(gap)
    return y


def _footer(slide, n, total):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), H - Inches(0.38), W, Inches(0.38))
    _fill(bar, FOOTER, line=LINE)
    _textbox(
        slide,
        Inches(0.45),
        H - Inches(0.36),
        Inches(10.2),
        Inches(0.32),
        "Internal briefing  ·  Saab  ·  KTH master's thesis  ·  Markus Johnson Swegmark",
        size=11,
        color=MUTED,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    _textbox(
        slide,
        W - Inches(1.4),
        H - Inches(0.36),
        Inches(1.05),
        Inches(0.32),
        f"{n} / {total}",
        size=11,
        color=MUTED,
        align=PP_ALIGN.RIGHT,
        anchor=MSO_ANCHOR.MIDDLE,
    )


def _notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def _bg(slide):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    _fill(shp, BG)
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.07))
    _fill(accent, CYAN)


def _card(slide, left, top, width, height, color=CARD):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    _fill(shp, color, line=LINE)
    shp.adjustments[0] = 0.08
    return shp


def _plot(slide, path, left, top, width, max_height=None):
    path = Path(path)
    im = Image.open(path)
    iw, ih = im.size
    height = width * ih / iw
    if max_height is not None and height > max_height:
        height = max_height
        width = height * iw / ih
    pad = Inches(0.08)
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left - pad,
        top - pad,
        width + 2 * pad,
        height + 2 * pad,
    )
    card.fill.solid()
    card.fill.fore_color.rgb = CARD
    card.line.color.rgb = LINE
    card.line.width = Pt(1.0)
    card.adjustments[0] = 0.04
    slide.shapes.add_picture(str(path), left, top, width, height)
    return width, height


def raster_pdf(src: Path, name: str, zoom: float = 2.6) -> Path:
    dst = CONV / name
    doc = pymupdf.open(src)
    pix = doc[0].get_pixmap(matrix=pymupdf.Matrix(zoom, zoom), alpha=False)
    pix.save(dst)
    return dst


P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"


def find_anim(stem: str) -> tuple[Path | None, Path | None]:
    """Return (unused_gif, mp4_or_poster). GIFs loop in PowerPoint, so we prefer MP4."""
    mp4 = ANIM / f"{stem}.mp4"
    mp4_hq = ANIM / "manim_media" / "videos" / "scenes" / "1080p30" / f"{stem}.mp4"
    poster = ANIM / f"{stem}_poster.png"
    video = mp4 if mp4.exists() else (mp4_hq if mp4_hq.exists() else None)
    return None, video if video else (poster if poster.exists() else None)


def _autoplay_hold_last_frame(slide) -> None:
    """Play once when the slide appears; freeze on the last frame (do not rewind)."""
    sld = slide._element
    for node in sld.xpath(".//p:cMediaNode"):
        node.set("showWhenStopped", "1")
        ctn = None
        for child in list(node):
            if child.tag.endswith("}cTn") or child.tag == "cTn":
                ctn = child
                break
        if ctn is not None:
            ctn.set("fill", "hold")
            ctn.set("restart", "never")
            for el in ctn.iter():
                if el.tag.endswith("}cond") and el.get("delay") == "indefinite":
                    el.set("delay", "0")


def add_movie_or_gif(slide, gif, video, left, top, width, height, poster=None):
    if video is not None and Path(video).suffix.lower() == ".mp4":
        poster_path = poster if poster and Path(poster).exists() else None
        kwargs = dict(left=left, top=top, width=width, height=height, mime_type="video/mp4")
        if poster_path:
            slide.shapes.add_movie(str(video), poster_frame_image=str(poster_path), **kwargs)
        else:
            slide.shapes.add_movie(str(video), **kwargs)
        _autoplay_hold_last_frame(slide)
        return "mp4"
    if gif is not None:
        slide.shapes.add_picture(str(gif), left, top, width, height)
        return "gif"
    if video is not None:
        slide.shapes.add_picture(str(video), left, top, width, height)
        return "png"
    _card(slide, left, top, width, height)
    _textbox(slide, left, top + height / 2 - Inches(0.2), width, Inches(0.4),
             "Animation not rendered yet", size=16, color=MUTED, align=PP_ALIGN.CENTER)
    return "missing"


class Deck:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = W
        self.prs.slide_height = H
        self.layout = self.prs.slide_layouts[6]
        self.slides = []

    def new(self):
        s = self.prs.slides.add_slide(self.layout)
        _bg(s)
        self.slides.append(s)
        return s

    def finish(self):
        total = len(self.slides)
        for i, s in enumerate(self.slides, 1):
            _footer(s, i, total)

    def save(self, path: Path):
        self.finish()
        self.prs.save(path)
        return path


def build() -> Path:
    hist_em = raster_pdf(FIG / "hist_em.pdf", "hist_em.png")

    d = Deck()

    # 1 title
    s = d.new()
    _textbox(s, Inches(0.7), Inches(1.45), Inches(12), Inches(0.35),
             "INTERNAL BRIEFING  ·  SAAB / KTH", size=14, color=CYAN, bold=True)
    _textbox(s, Inches(0.7), Inches(1.9), Inches(12), Inches(1.5),
             "Joint emitter and mode clustering\nfrom mixed PDW streams",
             size=34, color=INK, bold=True)
    _textbox(s, Inches(0.7), Inches(3.55), Inches(11.5), Inches(0.7),
             "What the data looks like, how the encoder works, and which scores actually moved.",
             size=18, color=MUTED)
    _textbox(s, Inches(0.7), Inches(4.45), Inches(12), Inches(0.4),
             "Problem   ·   Data   ·   Encoder   ·   RoPE / bias   ·   Takeaways",
             size=16, color=CYAN, bold=True)
    _textbox(s, Inches(0.7), Inches(5.5), Inches(11), Inches(0.7),
             "Markus Johnson Swegmark\nSupervisors: Alexander Karlsson, Elin Ohlman (Saab)  ·  Saikat Chatterjee (KTH)",
             size=14, color=INK)
    _notes(s, "Internal walkthrough, not a defence. Data and the two attention tricks first.")

    # 2 problem + two jobs
    s = d.new()
    _textbox(s, Inches(0.55), Inches(0.25), Inches(12), Inches(0.42), "One mixed stream, two partitions", size=26, color=INK, bold=True)
    _plot(s, GEN / "interleaved_stream.png", Inches(0.55), Inches(0.78), Inches(12.2), max_height=Inches(2.55))
    _card(s, Inches(0.55), Inches(3.55), Inches(6.0), Inches(3.15))
    _textbox(s, Inches(0.75), Inches(3.7), Inches(5.6), Inches(0.35), "Emitter  ·  deinterleaving", size=16, color=ORANGE, bold=True)
    _bullets(s, Inches(0.75), Inches(4.15), Inches(5.55), [
        "IDs unique only inside this recording.",
        "Unknown how many platforms are present.",
        "We recover a partition, not a name.",
    ], size=15, gap=0.55)
    _card(s, Inches(6.75), Inches(3.55), Inches(6.0), Inches(3.15))
    _textbox(s, Inches(6.95), Inches(3.7), Inches(5.6), Inches(0.35), "Mode  ·  operating schedule", size=16, color=CYAN, bold=True)
    _bullets(s, Inches(6.95), Inches(4.15), Inches(5.55), [
        "19-type catalogue, shared across recordings.",
        "Same mode can sit on several emitters.",
        "Inference still clusters — no catalogue softmax.",
    ], size=15, gap=0.55)
    _notes(s, "Classical chain deinterleaves first. We want both partitions from the mixed list. Answers at the end.")

    # 3 corpus + dense/sparse
    s = d.new()
    _textbox(s, Inches(0.55), Inches(0.22), Inches(12), Inches(0.4), "Dense lock-on versus sparse scan", size=24, color=INK, bold=True)
    stats = [("L = 2000", "pulses / window"), ("1 169", "test windows"), ("~3 / ~6", "emitters / modes"), ("19", "mode types")]
    for i, (v, lab) in enumerate(stats):
        x = Inches(0.55 + i * 3.2)
        _card(s, x, Inches(0.7), Inches(3.05), Inches(0.85))
        _textbox(s, x + Inches(0.1), Inches(0.74), Inches(2.85), Inches(0.4), v, size=16, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
        _textbox(s, x + Inches(0.1), Inches(1.12), Inches(2.85), Inches(0.32), lab, size=11, color=MUTED, align=PP_ALIGN.CENTER)
    _plot(s, FIG / "feat_w100_em.png", Inches(0.45), Inches(1.7), Inches(7.5), max_height=Inches(5.0))
    _bullets(s, Inches(8.2), Inches(1.85), Inches(4.7), [
        "Orange: lock-on, low PRI, fills the window.",
        "Purple: scanning. Only seen when the beam sweeps the receiver.",
        "Amplitude is the giveaway: lobes, not a line. Scan period 2–5 s.",
        "A density clusterer on RF/PW ignores the sparse class.",
        "3 pulses vs 1997 in the same window is common.",
    ], size=14, gap=0.72)
    _notes(s, "Walk the amplitude panel. Gaps between purple lobes are the rotation.")

    # 4 imbalance + crowded RF/PW
    s = d.new()
    _textbox(s, Inches(0.55), Inches(0.2), Inches(12), Inches(0.4), "Floods, and RF/PW that are not unique IDs", size=24, color=INK, bold=True)
    _plot(s, FIG / "feat_imbalance.png", Inches(0.4), Inches(0.7), Inches(6.1), max_height=Inches(4.5))
    _plot(s, FIG / "feat_maxemitters.png", Inches(6.75), Inches(0.7), Inches(6.1), max_height=Inches(4.5))
    _textbox(s, Inches(0.5), Inches(5.5), Inches(6.1), Inches(1.2),
             "Left: 3 pulses vs 1997. One lock-on, one long PRI or mid-scan. A 3-pulse train is a poor later mode input.",
             size=13, color=INK)
    _textbox(s, Inches(6.85), Inches(5.5), Inches(6.0), Inches(1.2),
             "Right: several platforms share RF and PW bands. Timing is the remaining cue. All-pairs attention is the bet.",
             size=13, color=INK)
    _notes(s, "Why attention rather than an O(L) state. Feature DBSCAN completeness 0.52 on emitters.")

    # 5 architecture
    s = d.new()
    _textbox(s, Inches(0.55), Inches(0.22), Inches(12), Inches(0.4), "Shared trunk, two heads, two positive sets", size=24, color=INK, bold=True)
    _plot(s, GEN / "y_architecture.png", Inches(0.4), Inches(0.75), Inches(7.7), max_height=Inches(5.85))
    _bullets(s, Inches(8.3), Inches(0.85), Inches(4.6), [
        "~16 M parameters. Width 256, 8 heads. L = 2000 in, L embeddings out.",
        "Self-attention is O(L²). That mix is deinterleaving.",
        "LSTM / Mamba are O(L), but pairwise “same train?” is a comparison, not a state. Not compared.",
        "Emitter loss: same recording-local ID. Mode loss: same global catalogue label.",
        "Joint training is the unweighted sum. DBSCAN at inference, not in the loss.",
    ], size=13, gap=0.78)
    _notes(s, "The two similarities can fight in the trunk. Do not claim a Mamba bake-off.")

    # 9 two mechanisms
    s = d.new()
    _textbox(s, Inches(0.55), Inches(0.22), Inches(12), Inches(0.4), "Two ways to put physics into attention", size=24, color=INK, bold=True)
    _card(s, Inches(0.55), Inches(0.9), Inches(6.0), Inches(5.7))
    _textbox(s, Inches(0.8), Inches(1.1), Inches(5.5), Inches(0.4), "RoPE-TOA", size=20, color=CYAN, bold=True)
    _bullets(s, Inches(0.8), Inches(1.65), Inches(5.5), [
        "Rotate Q and K by window-local ToA.",
        "The score is a dot product, so the enclosed angle matters.",
        "Only the relative rotation R(Δt) enters the product.",
        "No extra weights. Encoder ~2.5 ms.",
    ], size=15, gap=0.72)
    _card(s, Inches(6.8), Inches(0.9), Inches(6.0), Inches(5.7))
    _textbox(s, Inches(7.05), Inches(1.1), Inches(5.5), Inches(0.4), "Physical bias", size=20, color=ORANGE, bold=True)
    _bullets(s, Inches(7.05), Inches(1.65), Inches(5.5), [
        "Build B_ij = |p_i − p_j| from ToA and incidence.",
        "Scale by learned λ, subtract from the attention matrix A.",
        "Same emitter scores as RoPE-TOA.",
        "Encoder ~32 ms. Extra L×L maps in every head.",
    ], size=15, gap=0.72)
    _notes(s, "Next two slides are the animations. Let them play once.")

    # 10 RoPE
    s = d.new()
    _textbox(s, Inches(0.55), Inches(0.18), Inches(12), Inches(0.35), "RoPE-TOA  ·  the score is the rotated dot product", size=22, color=INK, bold=True)
    gif, vid = find_anim("RoPEToA")
    add_movie_or_gif(s, gif, vid, Inches(0.45), Inches(0.58), Inches(8.7), Inches(4.9), poster=ANIM / "RoPEToA_poster.png")
    _bullets(s, Inches(9.35), Inches(0.75), Inches(3.6), [
        "q · k = cos φ.",
        "Rotate each by its own ToA.",
        "New angle: φ + Δt · ω.",
        "Close in time → small extra twist, high score.",
        "Far in time → larger twist, lower score.",
        "Plays once, then holds the last frame.",
    ], size=13, gap=0.68)
    _notes(s, "Vanilla already has ToA as a channel. RoPE puts Δt into every query-key product.")

    # 11 bias
    s = d.new()
    _textbox(s, Inches(0.55), Inches(0.18), Inches(12), Inches(0.35), "Physical bias  ·  learned λ on delta maps, subtracted from A", size=22, color=INK, bold=True)
    gif, vid = find_anim("PhysicalBias")
    add_movie_or_gif(s, gif, vid, Inches(0.45), Inches(0.58), Inches(8.7), Inches(4.9), poster=ANIM / "PhysicalBias_poster.png")
    _bullets(s, Inches(9.35), Inches(0.75), Inches(3.6), [
        "|ΔToA| and |Δ incidence| from the inputs.",
        "λ_t, λ_u: learned scalars on ToA and incidence.",
        "Bias λB is subtracted from vanilla A.",
        "Same-direction, nearby pulses stay bright.",
        "Maps are L×L. That is the bill.",
    ], size=13, gap=0.72)
    _notes(s, "A cache skips rebuild, not the quadratic subtract. RoPE is the cheap twin.")

    # 12 vanilla + tail
    s = d.new()
    _textbox(s, Inches(0.55), Inches(0.2), Inches(12), Inches(0.4), "Vanilla already works on most windows", size=24, color=INK, bold=True)
    _plot(s, hist_em, Inches(0.4), Inches(0.7), Inches(7.35), max_height=Inches(5.9))
    _bullets(s, Inches(8.0), Inches(0.85), Inches(4.9), [
        "Feature DBSCAN emitter ARI 0.58. Completeness 0.52.",
        "Vanilla joint: 0.95 emitter, 0.99 mode.",
        "88% of windows ≥ 0.99. About 1 in 10 below 0.90.",
        "RoPE-TOA and Bias stay just below 1. None below 0.90.",
        "Stacking them is worse than Vanilla.",
    ], size=14, gap=0.82)
    _notes(s, "RQ1 is yes, with a leftover deinterleaving tail. Mode is already easier in the raw PDWs.")

    # 10 joint + bakeoff + cost
    s = d.new()
    _textbox(s, Inches(0.55), Inches(0.18), Inches(12), Inches(0.35), "Joint is a tax on Vanilla. Time in attention removes it.", size=20, color=INK, bold=True)
    _plot(s, GEN / "joint_tax.png", Inches(0.35), Inches(0.58), Inches(4.2), max_height=Inches(3.15))
    _plot(s, GEN / "emitter_ari.png", Inches(4.6), Inches(0.58), Inches(4.3), max_height=Inches(3.15))
    _plot(s, GEN / "encoder_time.png", Inches(9.0), Inches(0.58), Inches(3.95), max_height=Inches(3.15))
    _textbox(s, Inches(0.5), Inches(4.0), Inches(12.3), Inches(0.85),
             "Emitter-only Vanilla hits 0.999; joint drops to 0.948. RoPE-TOA and Bias close the tail (ARI 1.000) without that tax. "
             "RoPE-az 0.973. Stacking the two time devices scores below Vanilla. RoPE-TOA is +0.9 ms; Bias is ~32 ms encoder.",
             size=14, color=INK)
    _bullets(s, Inches(0.55), Inches(4.85), Inches(12.2), [
        "Both partitions from the mixed window → joint RoPE-TOA. Do not stack the bias.",
        "Deinterleaving only → emitter-only Vanilla is already at the ceiling. Cluster IDs die at the window edge.",
        "Synthetic. Held-out modulations and operational data are not scored.",
    ], size=14, gap=0.48)
    _notes(s, "RQ2 is no for deinterleaving alone. RQ3 is yes for time, no for stacking. Do not recommend Bias for production.")

    # 11 close
    s = d.new()
    _textbox(s, Inches(0.55), Inches(0.25), Inches(12), Inches(0.45), "In one breath", size=26, color=INK, bold=True)
    lines = [
        ("RQ1", "Yes. Vanilla recovers both partitions on most windows. Feature DBSCAN does not."),
        ("RQ2", "No, if you only need deinterleaving. Joint taxes emitters on Vanilla. RoPE-TOA does not pay that tax."),
        ("RQ3", "Yes for time. RoPE-TOA or the pairwise bias close the emitter tail. Stacking them does not. RoPE is the one you want."),
    ]
    for i, (tag, t) in enumerate(lines):
        y = Inches(1.0 + i * 1.45)
        _card(s, Inches(0.55), y, Inches(12.2), Inches(1.3))
        chip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y + Inches(0.4), Inches(1.2), Inches(0.48))
        _fill(chip, CYAN)
        chip.adjustments[0] = 0.2
        _textbox(s, Inches(0.8), y + Inches(0.44), Inches(1.2), Inches(0.42), tag, size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        _textbox(s, Inches(2.25), y + Inches(0.32), Inches(10.15), Inches(0.7), t, size=16, color=INK)
    _textbox(s, Inches(0.6), Inches(5.5), Inches(12.1), Inches(1.15),
             "The data is unbalanced because physics is: lock-on versus a 2–5 s scan, low PRI versus high PRI. "
             "Attention is expensive because associating non-adjacent pulses is an all-pairs question. Putting time into that product closed the last errors.",
             size=14, color=MUTED)
    _notes(s, "Stop. Questions. Mamba: not compared. Real data: not scored. Tracks: future work.")

    out = ROOT / "Saab_internal_briefing.pptx"
    d.save(out)
    return out


if __name__ == "__main__":
    path = build()
    print(path)
